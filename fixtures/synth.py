"""Generate deterministic fixture decks for local testing.

No LLM here — every piece of content is hand-authored so tests can assert
specific findings (e.g. "slide 1 should be flagged as non-action-title").
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1 and subtitle:
        slide.placeholders[1].text = subtitle


def _add_footer(
    slide,
    text: str,
    *,
    top_in: float = 6.95,
    left_in: float = 0.5,
    height_in: float = 0.3,
    width_in: float = 6.0,
    font_size_pt: int = 10,
) -> None:
    """Add a footer textbox at the bottom of a slide."""
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = tb.text_frame
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(font_size_pt)


def _add_content_slide(
    prs: Presentation,
    title: str,
    body: str,
    *,
    footer_text: str | None = None,
    footer_top_in: float = 6.95,
    footer_left_in: float = 0.5,
    footer_height_in: float = 0.3,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = body
    if footer_text:
        _add_footer(
            slide,
            footer_text,
            top_in=footer_top_in,
            left_in=footer_left_in,
            height_in=footer_height_in,
        )


def build_bad_deck(out_path: str | Path) -> Path:
    """A deck full of issues: descriptive titles, long paragraphs, no so-what,
    inconsistent footers (text drift + misaligned + missing on some slides),
    no overlap with filename keywords.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cover - descriptive title, subtitle without filename keywords
    _add_title_slide(prs, "Resultados 2024", "Presentación general")

    # Slide 2 - footer "Estudio 2024" (drifts from cover "Resultados 2024")
    _add_content_slide(
        prs,
        title="Análisis de ventas",
        body=(
            "Las ventas durante el año 2024 mostraron un comportamiento mixto en los distintos canales "
            "comerciales que opera la compañía, con un crecimiento muy fuerte en el canal digital pero "
            "con una contracción simultánea en el canal físico que se vio afectado por una baja del "
            "tráfico peatonal en los principales centros comerciales del país durante el segundo semestre."
        ),
        footer_text="Estudio 2024",
        footer_top_in=6.95,
        footer_height_in=0.3,
    )

    # Slide 3 - DUPLICATE title + footer in different caps style + misaligned vertically
    _add_content_slide(
        prs,
        title="Costos operativos",
        body=(
            "Los costos operativos crecieron 12% impactados principalmente por la inflación. "
            "Igualmente, se observa que la presión sobre los márgenes continuará en 2025 si no se "
            "implementan medidas concretas de eficiencia operativa transversal."
        ),
        footer_text="ESTUDIO 2024",
        footer_top_in=7.20,           # misaligned vertically vs slide 2's 6.95
        footer_height_in=0.42,
    )

    # Slide 4 - same title as slide 3 (duplicate), NO footer (coverage gap)
    _add_content_slide(
        prs,
        title="Costos operativos",
        body=(
            "El margen EBITDA cayó 3 puntos porcentuales. Esto se explica por el alza en costos "
            "logísticos y los mayores precios de materias primas durante todo el período. "
            "Adicionalmente, hubo presión competitiva en categorías clave."
        ),
        footer_text=None,
    )

    # Slide 5 - CONTENT WITHOUT TITLE: blank layout with substantial body text and no footer.
    blank_layout = prs.slide_layouts[6]
    s5 = prs.slides.add_slide(blank_layout)
    tb = s5.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = (
        "La compañía requiere revisar urgentemente su estrategia de canales para "
        "responder a la caída sostenida del tráfico físico. Las opciones evaluadas "
        "incluyen cierre selectivo de tiendas, reconversión a pickup points, y "
        "expansión del marketplace para sumar SKUs sin asumir inventario."
    )

    out = Path(out_path)
    prs.save(out)
    return out


def build_good_deck(out_path: str | Path) -> Path:
    """A clean deck: action titles, short bulleted body, consistent footers
    that repeat the deck title, filename keywords appear in cover.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    deck_title_footer = "Estrategia comercial 2025"

    # Cover - filename will be 'estrategia_comercial_2025.pptx', subtitle shares keywords
    _add_title_slide(
        prs,
        "Estrategia comercial 2025: capturar +20% de share digital",
        "Plan comercial — crecimiento del canal directo en 2025",
    )

    # Slide 2 - clear action title, body broken into short paragraphs, FOOTER
    _add_content_slide(
        prs,
        title="El canal digital creció 30% en 2024, superando al canal físico por segundo año consecutivo",
        body=(
            "Penetración digital pasó de 45% a 60% del total\n"
            "Nuevas categorías (electro, hogar) impulsaron 15pp del crecimiento\n"
            "Conversion rate subió de 1.8% a 2.5% tras rediseñar el checkout"
        ),
        footer_text=deck_title_footer,
    )

    # Slide 3 - cause then consequence, same consistent footer
    _add_content_slide(
        prs,
        title="Invertir USD 8M en logística de última milla aumentaría conversión en 4pp y margen en 2pp",
        body=(
            "La fricción de entrega es el principal driver de abandono de carrito\n"
            "Reducir tiempos a 24h en top-30 ciudades cubre 70% de la demanda\n"
            "Payback estimado: 14 meses según modelo financiero validado con CFO"
        ),
        footer_text=deck_title_footer,
    )

    # Slide 4 - storyline conclusion, same footer
    _add_content_slide(
        prs,
        title="Aprobar la inversión en Q1 desbloquea +20% de share digital para fin de 2025",
        body=(
            "Decisión requerida en próxima reunión de comité\n"
            "Riesgo bajo: tecnología y partners ya validados\n"
            "Win rápido: piloto en Buenos Aires y CDMX en 90 días"
        ),
        footer_text=deck_title_footer,
    )

    out = Path(out_path)
    prs.save(out)
    return out


def build_deck_with_image(out_path: str | Path) -> Path:
    """A small deck that includes an embedded picture, for testing visual extraction."""
    from io import BytesIO

    from PIL import Image, ImageDraw

    img = Image.new("RGB", (320, 200), color=(20, 80, 160))
    draw = ImageDraw.Draw(img)
    draw.rectangle([(20, 30), (300, 50)], fill=(255, 255, 255))
    draw.rectangle([(20, 60), (200, 80)], fill=(255, 200, 0))
    draw.rectangle([(20, 90), (260, 110)], fill=(0, 200, 100))
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, "Reporte con visuales", "Análisis Q4 con gráficos")

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Crecimiento por canal"
    slide.shapes.add_picture(buf, Inches(3), Inches(2), Inches(7), Inches(4))

    out = Path(out_path)
    prs.save(out)
    return out


BAD_DECK_FILENAME = "deck_problemas.pptx"
GOOD_DECK_FILENAME = "estrategia_comercial_2025.pptx"
IMAGE_DECK_FILENAME = "deck_con_imagen.pptx"
