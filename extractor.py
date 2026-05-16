from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def _emu_to_inches(value: int | None) -> float | None:
    if value is None:
        return None
    return round(Emu(value).inches, 2)


def _is_picture(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except (AttributeError, ValueError):
        return False


def _is_chart(shape) -> bool:
    try:
        return shape.has_chart
    except AttributeError:
        return False


def extract_deck(pptx_path: str | Path) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    slide_width = _emu_to_inches(prs.slide_width)
    slide_height = _emu_to_inches(prs.slide_height)

    slides: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_text = title_shape.text_frame.text.strip() if title_shape and title_shape.has_text_frame else None

        picture_count = sum(1 for s in slide.shapes if _is_picture(s))
        chart_count = sum(1 for s in slide.shapes if _is_chart(s))

        shapes_data: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            paragraphs = []
            shape_explicit_sizes_pt: list[float] = []
            for p in shape.text_frame.paragraphs:
                p_text = "".join(run.text for run in p.runs).strip()
                if not p_text:
                    continue
                run_sizes_pt: list[float] = []
                for run in p.runs:
                    if not (run.text or "").strip():
                        continue
                    sz = run.font.size
                    if sz is not None:
                        try:
                            run_sizes_pt.append(float(sz.pt))
                        except (AttributeError, ValueError):
                            pass
                min_size_pt = min(run_sizes_pt) if run_sizes_pt else None
                shape_explicit_sizes_pt.extend(run_sizes_pt)
                paragraphs.append({
                    "text": p_text,
                    "level": p.level,
                    "min_size_pt": min_size_pt,
                })

            info: dict[str, Any] = {
                "name": shape.name,
                "text": text,
                "paragraphs": paragraphs,
                "min_font_size_pt": (
                    min(shape_explicit_sizes_pt) if shape_explicit_sizes_pt else None
                ),
                "is_title": title_shape is not None and shape == title_shape,
                "top_in": _emu_to_inches(shape.top),
                "left_in": _emu_to_inches(shape.left),
                "width_in": _emu_to_inches(shape.width),
                "height_in": _emu_to_inches(shape.height),
            }
            if shape.is_placeholder:
                info["placeholder_type"] = str(shape.placeholder_format.type)
                info["placeholder_idx"] = shape.placeholder_format.idx

            shapes_data.append(info)

        notes_text = None
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip() or None

        slides.append(
            {
                "slide_number": i,
                "layout_name": slide.slide_layout.name,
                "title": title_text,
                "shapes": shapes_data,
                "notes": notes_text,
                "picture_count": picture_count,
                "chart_count": chart_count,
                "has_visuals": picture_count > 0 or chart_count > 0,
            }
        )

    return {
        "slide_width_in": slide_width,
        "slide_height_in": slide_height,
        "slide_count": len(slides),
        "slides": slides,
    }


def extract_images(pptx_path: str | Path) -> dict[int, list[tuple[bytes, str]]]:
    """Extract embedded picture bytes per slide.

    Returns a dict mapping slide_number -> list of (image_bytes, extension).
    Only slides with at least one picture appear in the dict. Charts are NOT
    rasterized — that would require a renderer dependency; they're left as
    placeholders in the visual count.
    """
    prs = Presentation(str(pptx_path))
    result: dict[int, list[tuple[bytes, str]]] = {}
    for i, slide in enumerate(prs.slides, start=1):
        images: list[tuple[bytes, str]] = []
        for shape in slide.shapes:
            if not _is_picture(shape):
                continue
            try:
                img = shape.image
                images.append((img.blob, img.ext))
            except (AttributeError, ValueError):
                continue
        if images:
            result[i] = images
    return result
