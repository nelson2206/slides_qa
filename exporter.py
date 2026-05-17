"""Holmes export pipeline — produce review artefacts from analysis results.

Two surfaces:
  - export_annotated_pptx(): take the original deck + a QA result, return a
    new .pptx with Holmes' findings injected into each slide's speaker notes
    plus a 'Holmes Review' summary slide appended at the end.
  - apply_quick_fixes(): take the original deck + a QA result + a list of fix
    selections, return a new .pptx with the chosen fixes baked into the slide
    content (sentence-case title rewrite, replace title with the LLM
    suggestion, replace footer text with the canonical text).

Both functions return bytes so the Streamlit UI can offer them as downloads
without writing to disk.
"""
from __future__ import annotations

import io
import re
from typing import Any, Literal

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from checks import NON_BODY_ROLES, _to_sentence_case


# ---------------------------------------------------------------------------
# Annotated export — Holmes' notes inside the speaker notes of each slide
# ---------------------------------------------------------------------------

_HOLMES_NOTES_HEADER = "\n\n— — — HOLMES REVIEW — — —\n"


def _format_slide_findings(finding: dict[str, Any]) -> str:
    """Build the speaker-notes block Holmes injects per slide."""
    n = finding["slide_number"]
    role = finding.get("role", "?")
    score = finding.get("score")
    sev = (finding.get("severity") or "nit").upper()
    skipped = finding.get("_skipped", False)

    score_str = f"{score}/10" if score is not None else "—"
    lines: list[str] = [
        f"Slide {n}  ·  Score {score_str}  ·  {sev}  ·  role={role}",
    ]
    if skipped:
        lines.append("(Skipped — slide estructural, sin evaluación semántica.)")
        return _HOLMES_NOTES_HEADER + "\n".join(lines) + "\n"

    summary = (finding.get("summary") or "").strip()
    if summary:
        lines.append(f"Resumen: {summary}")

    # Per-check findings
    def _block(label: str, status: str, suggestion: str | None = None,
               current: str | None = None) -> None:
        chunk = [f"\n[{label}]  {status}"]
        if current:
            chunk.append(f"   Actual: {current}")
        if suggestion:
            chunk.append(f"   Sugerencia: {suggestion}")
        lines.append("\n".join(chunk))

    at = finding.get("action_title") or {}
    verdict = at.get("is_action_title")
    status = "✓ OK" if verdict is True else "✗ FALLA" if verdict is False else "— no eval"
    _block("Action title", status,
           suggestion=at.get("suggestion"),
           current=at.get("current_title"))

    sw = finding.get("so_what") or {}
    present = sw.get("present")
    status = "✓ presente" if present is True else "✗ falta" if present is False else "— no eval"
    _block("So-what", status, suggestion=sw.get("suggestion"))

    cc = finding.get("cause_consequence") or {}
    ok = cc.get("ok")
    status = "✓ OK" if ok is True else "✗ invertido" if ok is False else "— no eval"
    _block("Causa → consecuencia", status)
    if cc.get("notes"):
        lines.append(f"   {cc['notes']}")

    tl = finding.get("text_length") or {}
    if not tl.get("ok"):
        _block(
            "Longitud de párrafos",
            tl.get("notes", "—"),
            suggestion=tl.get("suggestion"),
        )

    footer = finding.get("footer") or {}
    if not footer.get("exempt"):
        if not footer.get("present"):
            _block("Pie de página", "✗ ausente",
                   suggestion=(f'Agregar: "{footer.get("canonical_text")}"'
                               if footer.get("canonical_text") else None))
        elif footer.get("matches_canonical") is False:
            _block("Pie de página", "✗ texto distinto al canónico",
                   current=footer.get("current_footer"),
                   suggestion=(f'Reemplazar por: "{footer.get("canonical_text")}"'
                               if footer.get("canonical_text") else None))
        elif footer.get("aligned") is False:
            outlier = footer.get("alignment_outlier") or {}
            top = outlier.get("canonical_top_in")
            left = outlier.get("canonical_left_in")
            sugg = None
            if top is not None and left is not None:
                sugg = (
                    f"Mover a esquina inferior izquierda — top {top:.2f}″ · left {left:.2f}″ "
                    "(mediana de los pies del deck)."
                )
            _block("Pie de página", "✗ posición fuera del canónico", suggestion=sugg)

    tc = finding.get("title_case") or {}
    if tc.get("applicable") and not tc.get("ok"):
        violation = tc.get("case_violation", "casing")
        label = "MAYÚSCULAS" if violation == "all_caps" else "Title Case"
        _block(f"Casing del título · {label}", tc.get("notes", "—"),
               current=tc.get("title"), suggestion=tc.get("suggestion"))

    ff = finding.get("font_family") or {}
    if ff.get("applicable") and not ff.get("ok"):
        _block("Fuente brand", ff.get("notes", "—"),
               suggestion=ff.get("suggestion"))

    mfs = finding.get("min_font_size") or {}
    if mfs.get("applicable") and not mfs.get("ok"):
        _block(f"Tamaño mínimo ≥{int(mfs.get('min_required_pt', 9))}pt",
               mfs.get("notes", "—"), suggestion=mfs.get("suggestion"))

    td = finding.get("text_density") or {}
    if td.get("applicable") and not td.get("ok"):
        _block("Densidad de texto", td.get("notes", "—"),
               suggestion=td.get("suggestion"))

    return _HOLMES_NOTES_HEADER + "\n".join(lines) + "\n"


def _format_overview_summary(result: dict[str, Any]) -> str:
    """Multi-line summary text for the 'Holmes Review' appendix slide."""
    overview = result.get("deck_overview") or {}
    slides = result.get("slides") or []
    total = len(slides)
    sev_counts: dict[str, int] = {}
    for s in slides:
        sev = (s.get("severity") or "nit")
        sev_counts[sev] = sev_counts.get(sev, 0) + 1

    scores = [s.get("score") for s in slides if s.get("score") is not None]
    avg = sum(scores) / len(scores) if scores else 0

    lines = [
        "HOLMES — Auditoría de presentaciones",
        "",
        f"Score promedio: {avg:.1f}/10  ·  {total} slides analizados",
        (f"Critical: {sev_counts.get('critical', 0)}   "
         f"Warning: {sev_counts.get('warning', 0)}   "
         f"Nit: {sev_counts.get('nit', 0)}   "
         f"OK: {sev_counts.get('ok', 0)}"),
        "",
    ]

    coh = overview.get("storyline_coherent")
    if coh is not None:
        story = "✓ Coherente" if coh else "✗ Falta coherencia"
        lines.append(f"Storyline cross-slide: {story}")
        notes = overview.get("storyline_notes")
        if notes:
            lines.append(f"  {notes}")
        lines.append("")

    filename = overview.get("filename_subtitle_alignment")
    if filename:
        lines.append(f"Filename ↔ títulos: {filename}")
        lines.append("")

    cross = overview.get("cross_slide_issues") or []
    if cross:
        lines.append(f"Issues cross-slide ({len(cross)}):")
        for issue in cross:
            nums = ", ".join(str(s) for s in issue.get("slide_numbers", []))
            lines.append(f"  · Slides {nums}: {issue.get('issue', '')}")
        lines.append("")

    flagged = [s for s in slides
               if (s.get("severity") in ("critical", "warning"))
               and not s.get("_skipped")]
    if flagged:
        lines.append(f"Slides a revisar primero ({len(flagged)}):")
        for s in flagged[:20]:
            n = s["slide_number"]
            sev = (s.get("severity") or "").upper()
            title = ((s.get("action_title") or {}).get("current_title") or "")[:55]
            lines.append(f"  · #{n}  {sev}  {title}")
        if len(flagged) > 20:
            lines.append(f"  …y {len(flagged) - 20} más.")

    return "\n".join(lines)


def _set_notes_text(slide, text: str) -> None:
    """Replace the speaker notes text of a slide with `text`.

    NOTE: python-pptx's `text_frame.text = "..."` setter puts the entire
    string into a single run of a single paragraph — literal '\\n' chars
    survive in the run but PowerPoint does NOT render them as line breaks.
    We must split on \\n and emit one paragraph per line so PowerPoint
    actually shows the structured notes.
    """
    notes_slide = slide.notes_slide  # creates one if absent
    tf = notes_slide.notes_text_frame

    # Clear any existing content so we don't accumulate paragraphs from
    # earlier calls or from the original deck's notes.
    tf.clear()

    lines = text.split("\n") if text else [""]
    if not lines:
        return

    # The first paragraph already exists after .clear() — reuse it.
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = lines[0]
    else:
        first_para.text = lines[0]

    # Each subsequent line becomes its own paragraph so the line breaks
    # render correctly in the PowerPoint notes pane.
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line


def _append_summary_slide(prs: Presentation, summary_text: str) -> None:
    """Add a 'Holmes Review' summary slide at the end of the deck."""
    # Pick a layout that gives us a title + content placeholder if possible
    layouts = prs.slide_layouts
    # 5 is typically "Title Only" — safer than relying on body placeholders that
    # vary by template. We add a manual text box for the body.
    layout_idx = 5 if len(layouts) > 5 else 0
    slide = prs.slides.add_slide(layouts[layout_idx])

    if slide.shapes.title is not None:
        slide.shapes.title.text = "Holmes Review"

    # Add a generous text box for the body summary
    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(prs.slide_width / 914400 - 1.0)  # full width minus margins
    height = Inches(prs.slide_height / 914400 - 1.8)
    body = slide.shapes.add_textbox(left, top, width, height)
    tf = body.text_frame
    tf.word_wrap = True
    tf.text = summary_text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(11)


def export_annotated_pptx(deck_path: str, result: dict[str, Any]) -> bytes:
    """Build an annotated copy of the deck and return it as bytes.

    Strategy:
      1. Inject Holmes' per-slide findings into the speaker notes of each
         existing slide (preserves any prior notes — Holmes' block is
         appended with a clear separator).
      2. Append a final 'Holmes Review' summary slide with the deck-level
         overview + a triage list of slides flagged critical/warning.

    The original visual content is untouched.
    """
    prs = Presentation(deck_path)
    findings_by_num = {f["slide_number"]: f for f in (result.get("slides") or [])}

    for i, slide in enumerate(prs.slides, start=1):
        finding = findings_by_num.get(i)
        if not finding:
            continue
        block = _format_slide_findings(finding)
        # Preserve original notes if any
        notes_slide = slide.notes_slide
        existing = (notes_slide.notes_text_frame.text or "").strip()
        new_text = (existing + block) if existing else block.lstrip("\n")
        _set_notes_text(slide, new_text)

    # Append summary slide
    summary = _format_overview_summary(result)
    _append_summary_slide(prs, summary)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Auto-fix engine — apply user-selected suggestions to the deck
# ---------------------------------------------------------------------------

FixId = Literal[
    "sentence_case_title",
    "apply_llm_action_title",
    "canonical_footer_text",
    "add_canonical_footer",
    "move_footer_to_canonical_position",
    "enforce_min_font_size",
    "enforce_brand_font",
    "bullet_long_paragraph",
]

# Sentence splitter: cuts on '.', '!', '?' followed by whitespace + capital
# letter (covering Spanish accents). Won't split on abbreviations because they
# rarely end in capital-letter-next-word.
_SENTENCE_SPLIT_RE = re.compile(r"(?<=[.!?])\s+(?=[A-ZÁÉÍÓÚÑ])")


def _split_into_sentences(text: str) -> list[str]:
    """Split text into sentences (Spanish-aware). Skips very short fragments."""
    raw = _SENTENCE_SPLIT_RE.split(text.strip())
    return [s.strip() for s in raw if len(s.strip()) >= 8]


def available_fixes_for_slide(finding: dict[str, Any]) -> list[dict[str, Any]]:
    """Return the list of auto-applicable fixes for a single slide finding.

    Each fix dict: {id, slide_number, label, preview_before, preview_after, …extra}
    """
    if finding.get("_skipped"):
        return []
    fixes: list[dict[str, Any]] = []
    n = finding["slide_number"]

    # 1. Casing fix → convert ALL CAPS or Title Case to sentence case
    tc = finding.get("title_case") or {}
    if tc.get("applicable") and not tc.get("ok"):
        title = tc.get("title", "")
        suggested = _to_sentence_case(title)
        if suggested and suggested != title:
            fixes.append({
                "id": "sentence_case_title",
                "slide_number": n,
                "label": "Convertir título a sentence case",
                "preview_before": title,
                "preview_after": suggested,
            })

    # 2. Apply LLM action-title suggestion (only if explicitly suggested)
    at = finding.get("action_title") or {}
    if at.get("is_action_title") is False and at.get("suggestion"):
        current = at.get("current_title") or ""
        suggestion = at["suggestion"]
        if suggestion and suggestion != current:
            fixes.append({
                "id": "apply_llm_action_title",
                "slide_number": n,
                "label": "Reemplazar título con sugerencia de Holmes",
                "preview_before": current,
                "preview_after": suggestion,
            })

    # 3. FOOTER fixes — three distinct cases
    footer = finding.get("footer") or {}
    canonical = footer.get("canonical_text")
    exempt = footer.get("exempt", False)

    # 3a. Footer present but text differs → replace text
    if (
        canonical
        and not exempt
        and footer.get("present")
        and footer.get("matches_canonical") is False
    ):
        fixes.append({
            "id": "canonical_footer_text",
            "slide_number": n,
            "label": "Reemplazar pie de página con texto canónico",
            "preview_before": footer.get("current_footer") or "",
            "preview_after": canonical,
        })

    # 3b. Footer present but in wrong position → move to canonical top/left
    outlier = footer.get("alignment_outlier") or {}
    if (
        not exempt
        and footer.get("present")
        and outlier
        and outlier.get("canonical_top_in") is not None
        and outlier.get("canonical_left_in") is not None
    ):
        ct = outlier["canonical_top_in"]
        cl = outlier["canonical_left_in"]
        cur_t = outlier.get("current_top_in", 0)
        cur_l = outlier.get("current_left_in", 0)
        fixes.append({
            "id": "move_footer_to_canonical_position",
            "slide_number": n,
            "label": "Mover pie de página a esquina inferior izquierda",
            "preview_before": f"top {cur_t:.2f}″ · left {cur_l:.2f}″",
            "preview_after": f"top {ct:.2f}″ · left {cl:.2f}″",
            "canonical_top_in": ct,
            "canonical_left_in": cl,
        })

    # 3c. Footer absent → add it (only if canonical text + position known)
    if (
        canonical
        and not exempt
        and not footer.get("present")
        and footer.get("canonical_top_in") is not None
        and footer.get("canonical_left_in") is not None
    ):
        ct = footer["canonical_top_in"]
        cl = footer["canonical_left_in"]
        ch = footer.get("canonical_height_in") or 0.3
        fixes.append({
            "id": "add_canonical_footer",
            "slide_number": n,
            "label": "Agregar pie de página canónico",
            "preview_before": "(sin pie de página)",
            "preview_after": f'"{canonical}" en top {ct:.2f}″ · left {cl:.2f}″',
            "canonical_text": canonical,
            "canonical_top_in": ct,
            "canonical_left_in": cl,
            "canonical_height_in": ch,
        })

    # 4. Min font size — bump runs below 9pt up to 9pt
    mfs = finding.get("min_font_size") or {}
    if mfs.get("applicable") and not mfs.get("ok"):
        min_required = mfs.get("min_required_pt", 9.0)
        smallest = mfs.get("smallest_pt", 0)
        violations = mfs.get("violations") or []
        fixes.append({
            "id": "enforce_min_font_size",
            "slide_number": n,
            "label": f"Subir fuentes < {int(min_required)}pt al mínimo",
            "preview_before": f"{len(violations)} shape(s) · más chico {smallest:.1f}pt",
            "preview_after": f"todos los runs explícitos ≥ {int(min_required)}pt",
            "min_pt": min_required,
        })

    # 5. Brand font — replace non-ForFuture font names
    ff = finding.get("font_family") or {}
    if ff.get("applicable") and not ff.get("ok"):
        non_brand = ff.get("non_brand") or []
        non_brand_names = sorted({v["font"] for v in non_brand})
        if non_brand_names:
            fixes.append({
                "id": "enforce_brand_font",
                "slide_number": n,
                "label": "Reemplazar fuentes off-brand con ForFuture Sans",
                "preview_before": ", ".join(non_brand_names[:3])
                                  + ("…" if len(non_brand_names) > 3 else ""),
                "preview_after": "ForFuture Sans",
                "brand_font": "ForFuture Sans",
            })

    # 6. Bullet a long single paragraph
    tl = finding.get("text_length") or {}
    long_paras = tl.get("long_paragraphs") or []
    # Only offer when there's exactly one long paragraph (clean case) and it
    # has at least 2 sentences worth splitting.
    if len(long_paras) == 1:
        # The long_paragraphs entries are pre-formatted strings; we need the
        # raw shape name + snippet to find the actual shape in apply step.
        # Encode the index in the fix payload.
        raw_long = long_paras[0]
        # Extract shape name from "(ShapeName, ~N líneas, M palabras) …snippet"
        m = re.match(r"\(([^,]+),", raw_long)
        if m:
            shape_name = m.group(1).strip()
            fixes.append({
                "id": "bullet_long_paragraph",
                "slide_number": n,
                "label": "Dividir párrafo largo en bullets (por oración)",
                "preview_before": raw_long[:140] + ("…" if len(raw_long) > 140 else ""),
                "preview_after": "1 párrafo → N bullets cortos (split por oración)",
                "shape_name": shape_name,
            })

    return fixes


def _set_title_text(slide, new_text: str) -> bool:
    """Replace the title shape text. Returns True if a title was found."""
    title_shape = slide.shapes.title
    if title_shape is None or not title_shape.has_text_frame:
        return False
    tf = title_shape.text_frame
    # Preserve formatting of the first run when possible
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = new_text
        # Clear extra runs in first paragraph
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
        # Clear extra paragraphs
        for para in list(tf.paragraphs[1:]):
            for run in para.runs:
                run.text = ""
    else:
        tf.text = new_text
    return True


def _replace_footer_text(slide, current_text: str, new_text: str,
                        slide_height_in: float | None) -> bool:
    """Find the footer shape (bottom-left strip) and replace its text."""
    if slide_height_in is None:
        return False
    bottom_threshold = slide_height_in * 0.82
    # Find candidate footer shapes — same heuristic as checks.check_footer
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text:
            continue
        try:
            top_in = shape.top / 914400 if shape.top is not None else None
            height_in = shape.height / 914400 if shape.height is not None else None
            left_in = shape.left / 914400 if shape.left is not None else None
        except (TypeError, AttributeError):
            continue
        if top_in is None or height_in is None or left_in is None:
            continue
        if top_in < bottom_threshold or height_in >= 0.7:
            continue
        candidates.append((left_in, shape, text))

    if not candidates:
        return False
    candidates.sort(key=lambda c: c[0])  # leftmost first
    _, footer_shape, _ = candidates[0]
    tf = footer_shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = new_text
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
        for para in list(tf.paragraphs[1:]):
            for run in para.runs:
                run.text = ""
    else:
        tf.text = new_text
    return True


def _find_footer_shape(slide, slide_height_in: float | None):
    """Locate the footer shape in a slide (bottom-left strip)."""
    if slide_height_in is None:
        return None
    bottom_threshold = slide_height_in * 0.82
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text:
            continue
        try:
            top_in = shape.top / 914400 if shape.top is not None else None
            height_in = shape.height / 914400 if shape.height is not None else None
            left_in = shape.left / 914400 if shape.left is not None else None
        except (TypeError, AttributeError):
            continue
        if top_in is None or height_in is None or left_in is None:
            continue
        if top_in < bottom_threshold or height_in >= 0.7:
            continue
        candidates.append((left_in, shape))
    if not candidates:
        return None
    candidates.sort(key=lambda c: c[0])
    return candidates[0][1]


def _move_footer(slide, slide_height_in: float | None,
                 new_top_in: float, new_left_in: float) -> bool:
    """Find the footer and move it to (new_top_in, new_left_in) inches."""
    footer = _find_footer_shape(slide, slide_height_in)
    if footer is None:
        return False
    footer.top = Emu(int(new_top_in * 914400))
    footer.left = Emu(int(new_left_in * 914400))
    return True


def _add_footer(slide, text: str, top_in: float, left_in: float,
                height_in: float = 0.3, slide_width_in: float | None = None) -> bool:
    """Add a new textbox containing the canonical footer at the given position."""
    # Width: from left edge to ~60% of the slide width as a sensible default.
    if slide_width_in is None:
        width_in = 5.0
    else:
        width_in = max(2.0, min(slide_width_in - left_in - 0.3, slide_width_in * 0.6))
    tb = slide.shapes.add_textbox(
        Emu(int(left_in * 914400)),
        Emu(int(top_in * 914400)),
        Emu(int(width_in * 914400)),
        Emu(int(height_in * 914400)),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    # Style the run so it visually reads as a footer
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(9)
            run.font.name = "ForFuture Sans"
    return True


def _enforce_min_font(slide, min_pt: float) -> int:
    """Bump every run with explicit size < min_pt up to min_pt. Returns count."""
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size
                if sz is None:
                    continue
                try:
                    current_pt = float(sz.pt)
                except (AttributeError, ValueError):
                    continue
                if current_pt < min_pt - 0.01:
                    run.font.size = Pt(min_pt)
                    count += 1
    return count


_BRAND_FONT_NAME = "ForFuture Sans"
_BRAND_ACCEPTABLE = {
    "forfuture sans", "forfuturesans", "forfuture",
}


def _enforce_brand_font(slide) -> int:
    """Replace every run.font.name that's not in the brand-acceptable set
    with the canonical brand font. Returns count of runs updated."""
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                current = (run.font.name or "").strip().lower()
                if not current:
                    continue
                # Strip common style suffixes the way checks._normalize_font does
                norm = current.replace("-", " ").replace("_", " ")
                norm = " ".join(norm.split())
                if norm.replace(" ", "") == "forfuturesans":
                    continue
                if norm in _BRAND_ACCEPTABLE:
                    continue
                run.font.name = _BRAND_FONT_NAME
                count += 1
    return count


def _split_long_paragraph(slide, shape_name: str) -> tuple[bool, int]:
    """Find the named shape, take its (single) long paragraph and split it
    into one paragraph per sentence. Returns (ok, num_bullets)."""
    target = None
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            target = shape
            break
    if target is None:
        return False, 0
    tf = target.text_frame
    # Combine all paragraph texts (in case there are blank paragraphs)
    full_text = "\n".join(p.text for p in tf.paragraphs if p.text).strip()
    if not full_text:
        return False, 0
    sentences = _split_into_sentences(full_text)
    if len(sentences) < 2:
        return False, 0
    # Clear the text frame and write one paragraph per sentence
    tf.clear()
    first_para = tf.paragraphs[0]
    first_para.text = sentences[0]
    for s in sentences[1:]:
        p = tf.add_paragraph()
        p.text = s
    return True, len(sentences)


def apply_quick_fixes(
    deck_path: str,
    result: dict[str, Any],
    fixes_to_apply: list[dict[str, Any]],
) -> tuple[bytes, dict[str, Any]]:
    """Apply the selected fixes to a copy of the deck and return (bytes, report).

    `fixes_to_apply` is a list of fix dicts (id, slide_number, preview_after,
    ...) — the same shape returned by `available_fixes_for_slide`.

    Report has keys: applied (list), failed (list), counts (dict).
    """
    prs = Presentation(deck_path)
    slide_height_in = prs.slide_height / 914400 if prs.slide_height else None
    slide_width_in = prs.slide_width / 914400 if prs.slide_width else None
    slides_by_num = {i + 1: s for i, s in enumerate(prs.slides)}

    applied: list[dict[str, Any]] = []
    failed: list[dict[str, Any]] = []

    for fix in fixes_to_apply:
        n = fix["slide_number"]
        slide = slides_by_num.get(n)
        if slide is None:
            failed.append({**fix, "reason": "Slide no encontrada"})
            continue

        fix_id = fix["id"]
        try:
            if fix_id in ("sentence_case_title", "apply_llm_action_title"):
                ok = _set_title_text(slide, fix["preview_after"])
                if not ok:
                    failed.append({**fix, "reason": "Slide sin título editable"})
                    continue

            elif fix_id == "canonical_footer_text":
                ok = _replace_footer_text(
                    slide, fix["preview_before"], fix["preview_after"],
                    slide_height_in,
                )
                if not ok:
                    failed.append({**fix, "reason": "Pie de página no localizable"})
                    continue

            elif fix_id == "move_footer_to_canonical_position":
                ok = _move_footer(
                    slide, slide_height_in,
                    new_top_in=fix["canonical_top_in"],
                    new_left_in=fix["canonical_left_in"],
                )
                if not ok:
                    failed.append({**fix, "reason": "Pie de página no localizable para mover"})
                    continue

            elif fix_id == "add_canonical_footer":
                ok = _add_footer(
                    slide, fix["canonical_text"],
                    top_in=fix["canonical_top_in"],
                    left_in=fix["canonical_left_in"],
                    height_in=fix.get("canonical_height_in", 0.3),
                    slide_width_in=slide_width_in,
                )
                if not ok:
                    failed.append({**fix, "reason": "No se pudo crear el textbox"})
                    continue

            elif fix_id == "enforce_min_font_size":
                bumped = _enforce_min_font(slide, fix.get("min_pt", 9.0))
                if bumped == 0:
                    failed.append({**fix, "reason": "Sin runs de fuente explícita para corregir"})
                    continue

            elif fix_id == "enforce_brand_font":
                replaced = _enforce_brand_font(slide)
                if replaced == 0:
                    failed.append({**fix, "reason": "Sin runs con fuente off-brand explícita"})
                    continue

            elif fix_id == "bullet_long_paragraph":
                ok, bullets = _split_long_paragraph(slide, fix["shape_name"])
                if not ok:
                    failed.append({**fix, "reason": "No se pudo dividir el párrafo (<2 oraciones)"})
                    continue
                fix = {**fix, "bullets_created": bullets}

            else:
                failed.append({**fix, "reason": f"Fix id desconocido: {fix_id}"})
                continue

            applied.append(fix)
        except Exception as e:  # noqa: BLE001
            failed.append({**fix, "reason": str(e)})

    buf = io.BytesIO()
    prs.save(buf)

    report = {
        "applied": applied,
        "failed": failed,
        "counts": {
            "applied": len(applied),
            "failed": len(failed),
            "total_requested": len(fixes_to_apply),
        },
    }
    return buf.getvalue(), report
