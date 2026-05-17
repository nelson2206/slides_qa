from __future__ import annotations

import io

import pytest
from pptx import Presentation

from exporter import (
    apply_quick_fixes,
    available_fixes_for_slide,
    export_annotated_pptx,
)


# ---------- annotated export ----------

def _fake_result(slide_numbers: list[int]) -> dict:
    """Minimal result dict shape for export tests."""
    slides = []
    for n in slide_numbers:
        slides.append({
            "slide_number": n,
            "role": "content_with_title",
            "score": 6,
            "severity": "warning",
            "summary": f"Slide {n} test summary.",
            "_skipped": False,
            "action_title": {
                "is_action_title": False,
                "current_title": f"Title {n}",
                "notes": "Solo descriptivo.",
                "suggestion": f"Las ventas {n} cayeron 18%.",
            },
            "so_what": {
                "present": False,
                "notes": "Falta implicación.",
                "suggestion": "Agregar qué decisión habilita.",
            },
            "cause_consequence": {"ok": True, "notes": "OK"},
            "text_length": {"ok": True, "long_paragraphs": [], "notes": "OK", "suggestion": None},
            "footer": {
                "present": True,
                "current_footer": "Footer wrong",
                "aligned": True,
                "matches_canonical": False,
                "canonical_text": "Minsait · 2026",
                "exempt": False,
                "alignment_outlier": None,
                "notes": "—",
            },
            "title_case": {
                "applicable": True,
                "ok": False,
                "title": f"TITLE {n}",
                "case_violation": "all_caps",
                "notes": "MAYÚSCULAS",
                "suggestion": f'Reescribir en sentence case: "Title {n}".',
            },
            "font_family": {"applicable": True, "ok": True, "notes": "OK"},
            "min_font_size": {"applicable": True, "ok": True, "notes": "OK"},
            "text_density": {"applicable": True, "ok": True, "notes": "OK"},
        })
    return {
        "mode": "local",
        "slides": slides,
        "deck_overview": {
            "storyline_coherent": False,
            "storyline_notes": "Falta governing thought.",
            "filename_subtitle_alignment": "Filename no matchea.",
            "cross_slide_issues": [
                {"slide_numbers": [1, 2], "issue": "S1 y S2 invertidos."},
            ],
        },
    }


def test_export_annotated_pptx_writes_notes_and_summary_slide(bad_deck_path):
    deck = Presentation(str(bad_deck_path))
    original_slide_count = len(deck.slides)
    result = _fake_result(list(range(1, original_slide_count + 1)))

    out_bytes = export_annotated_pptx(str(bad_deck_path), result)
    assert len(out_bytes) > 0

    out_deck = Presentation(io.BytesIO(out_bytes))
    # +1 because of the summary slide appended at the end
    assert len(out_deck.slides) == original_slide_count + 1

    # Speaker notes on every existing slide should contain Holmes' block
    for i, slide in enumerate(out_deck.slides):
        if i == original_slide_count:
            break  # skip the summary slide we appended
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "HOLMES REVIEW" in notes_text, f"Slide {i+1} missing Holmes notes"
        # And the slide-specific number should be in this slide's notes,
        # not just slide 1's. Regression guard for the python-pptx text-setter
        # bug where all notes ended up as one blob in slide 1.
        assert f"Slide {i + 1}" in notes_text, (
            f"Slide {i+1} notes don't reference slide number "
            f"(actual: {notes_text[:200]!r})"
        )

    # Summary slide title
    summary_slide = out_deck.slides[-1]
    if summary_slide.shapes.title is not None:
        assert summary_slide.shapes.title.text == "Holmes Review"


def test_export_annotated_pptx_creates_multiple_paragraphs(bad_deck_path):
    """python-pptx's text_frame.text setter only sets the first paragraph —
    multi-line Holmes notes must be emitted as separate paragraphs so
    PowerPoint actually renders the line breaks."""
    deck = Presentation(str(bad_deck_path))
    original_slide_count = len(deck.slides)
    result = _fake_result(list(range(1, original_slide_count + 1)))

    out_bytes = export_annotated_pptx(str(bad_deck_path), result)
    out_deck = Presentation(io.BytesIO(out_bytes))

    for i, slide in enumerate(out_deck.slides):
        if i == original_slide_count:
            break
        tf = slide.notes_slide.notes_text_frame
        # Each slide's notes should have many paragraphs (header line, blank
        # lines, score line, summary, per-check blocks). Bare minimum > 3.
        assert len(tf.paragraphs) > 3, (
            f"Slide {i+1} notes have only {len(tf.paragraphs)} paragraph(s); "
            "python-pptx may have collapsed them into a single blob."
        )


def test_export_annotated_pptx_preserves_existing_notes(bad_deck_path):
    """Existing speaker notes should survive — Holmes' block is appended."""
    deck = Presentation(str(bad_deck_path))
    deck.slides[0].notes_slide.notes_text_frame.text = "Notas previas del consultor"
    modified_path = str(bad_deck_path) + ".modified.pptx"
    deck.save(modified_path)

    result = _fake_result([1])
    out_bytes = export_annotated_pptx(modified_path, result)
    out_deck = Presentation(io.BytesIO(out_bytes))
    notes = out_deck.slides[0].notes_slide.notes_text_frame.text
    assert "Notas previas del consultor" in notes
    assert "HOLMES REVIEW" in notes


# ---------- auto-fix engine ----------

def test_available_fixes_for_slide_lists_actionable_fixes():
    finding = {
        "slide_number": 3,
        "_skipped": False,
        "title_case": {
            "applicable": True,
            "ok": False,
            "title": "ANÁLISIS DE VENTAS",
            "case_violation": "all_caps",
        },
        "action_title": {
            "is_action_title": False,
            "current_title": "Análisis de ventas",
            "suggestion": "Las ventas cayeron 18% en Q3.",
        },
        "footer": {
            "present": True,
            "current_footer": "Wrong",
            "matches_canonical": False,
            "canonical_text": "Minsait · 2026",
            "exempt": False,
        },
    }
    fixes = available_fixes_for_slide(finding)
    ids = [f["id"] for f in fixes]
    assert "sentence_case_title" in ids
    assert "apply_llm_action_title" in ids
    assert "canonical_footer_text" in ids
    assert all(f["slide_number"] == 3 for f in fixes)


def test_available_fixes_for_slide_skips_skipped_slides():
    finding = {"slide_number": 1, "_skipped": True, "title_case": {"applicable": True, "ok": False}}
    assert available_fixes_for_slide(finding) == []


def test_available_fixes_for_slide_skips_when_no_diff():
    """If sentence-case version equals the original, don't offer the fix."""
    finding = {
        "slide_number": 2,
        "_skipped": False,
        "title_case": {
            "applicable": True,
            "ok": False,
            "title": "Already sentence case",
            "case_violation": "title_case",
        },
        "action_title": {},
        "footer": {},
    }
    fixes = available_fixes_for_slide(finding)
    # "Already sentence case" → _to_sentence_case → "Already sentence case" → no diff
    # (capitalized first, rest lowered; "already" → "Already" + "sentence" + "case" = "Already sentence case")
    ids = [f["id"] for f in fixes]
    assert "sentence_case_title" not in ids


def test_apply_quick_fixes_rewrites_title_text(good_deck_path):
    """Apply a sentence_case_title fix and verify the title text changed."""
    fix = {
        "id": "sentence_case_title",
        "slide_number": 2,
        "label": "Convertir título a sentence case",
        "preview_before": "ORIGINAL",
        "preview_after": "Sentence case title",
    }
    out_bytes, report = apply_quick_fixes(str(good_deck_path), {"slides": []}, [fix])
    assert report["counts"]["applied"] == 1
    assert report["counts"]["failed"] == 0

    out_deck = Presentation(io.BytesIO(out_bytes))
    slide_2 = out_deck.slides[1]
    if slide_2.shapes.title is not None:
        assert slide_2.shapes.title.text == "Sentence case title"


def test_apply_quick_fixes_reports_failures_gracefully(good_deck_path):
    """Asking to fix a slide that doesn't exist should fail cleanly, not crash."""
    fix = {
        "id": "sentence_case_title",
        "slide_number": 999,
        "preview_before": "X",
        "preview_after": "Y",
    }
    out_bytes, report = apply_quick_fixes(str(good_deck_path), {"slides": []}, [fix])
    assert report["counts"]["applied"] == 0
    assert report["counts"]["failed"] == 1
    assert "no encontrada" in report["failed"][0]["reason"].lower()
    # PPTX should still be returned (unmodified)
    out_deck = Presentation(io.BytesIO(out_bytes))
    assert len(out_deck.slides) > 0


# ---------- new fixes: footer move / add ----------

def test_available_fixes_offers_move_footer_when_outlier(bad_deck_path):
    """Slide whose footer is in the wrong position should get a move fix."""
    finding = {
        "slide_number": 2,
        "_skipped": False,
        "footer": {
            "present": True,
            "exempt": False,
            "alignment_outlier": {
                "current_top_in": 6.5, "current_left_in": 2.0,
                "canonical_top_in": 7.0, "canonical_left_in": 0.5,
                "issues": ["top=6.50in vs canónico 7.00in"],
            },
        },
        "action_title": {}, "title_case": {}, "min_font_size": {}, "font_family": {},
        "text_length": {},
    }
    fixes = available_fixes_for_slide(finding)
    move_fixes = [f for f in fixes if f["id"] == "move_footer_to_canonical_position"]
    assert len(move_fixes) == 1
    assert move_fixes[0]["canonical_top_in"] == 7.0
    assert move_fixes[0]["canonical_left_in"] == 0.5


def test_available_fixes_offers_add_footer_when_absent_with_canonical(bad_deck_path):
    """Slide without footer + deck has canonical text+position should get add fix."""
    finding = {
        "slide_number": 5,
        "_skipped": False,
        "footer": {
            "present": False,
            "exempt": False,
            "canonical_text": "Minsait · Project Acceleration",
            "canonical_top_in": 6.95,
            "canonical_left_in": 0.5,
            "canonical_height_in": 0.3,
        },
        "action_title": {}, "title_case": {}, "min_font_size": {}, "font_family": {},
        "text_length": {},
    }
    fixes = available_fixes_for_slide(finding)
    add_fixes = [f for f in fixes if f["id"] == "add_canonical_footer"]
    assert len(add_fixes) == 1
    assert add_fixes[0]["canonical_text"] == "Minsait · Project Acceleration"


def test_apply_move_footer_changes_position(good_deck_path):
    """move_footer_to_canonical_position should physically relocate the footer shape."""
    # First, find a slide in the good deck that already has a footer
    prs_check = Presentation(str(good_deck_path))
    target_slide_n = None
    for i, s in enumerate(prs_check.slides, start=1):
        # Quick footer detection: a small bottom shape with text
        for shape in s.shapes:
            if not shape.has_text_frame: continue
            if shape == s.shapes.title: continue
            text = (shape.text_frame.text or "").strip()
            if not text: continue
            try:
                top_in = shape.top / 914400 if shape.top else None
                h_in = shape.height / 914400 if shape.height else None
            except (TypeError, AttributeError):
                continue
            if top_in and h_in and top_in > 5.5 and h_in < 0.7:
                target_slide_n = i
                break
        if target_slide_n: break
    if target_slide_n is None:
        pytest.skip("Good deck has no detectable footer to move")

    fix = {
        "id": "move_footer_to_canonical_position",
        "slide_number": target_slide_n,
        "canonical_top_in": 7.0,
        "canonical_left_in": 0.3,
        "preview_before": "x", "preview_after": "y",
    }
    out_bytes, report = apply_quick_fixes(str(good_deck_path), {"slides": []}, [fix])
    assert report["counts"]["applied"] == 1
    # Verify the footer moved
    out_deck = Presentation(io.BytesIO(out_bytes))
    target = out_deck.slides[target_slide_n - 1]
    found_at_canonical = False
    for shape in target.shapes:
        if not shape.has_text_frame: continue
        if shape == target.shapes.title: continue
        if not (shape.text_frame.text or "").strip(): continue
        try:
            top_in = shape.top / 914400 if shape.top else None
            left_in = shape.left / 914400 if shape.left else None
        except (TypeError, AttributeError):
            continue
        if top_in and left_in and abs(top_in - 7.0) < 0.05 and abs(left_in - 0.3) < 0.05:
            found_at_canonical = True
            break
    assert found_at_canonical, "Footer was not relocated to canonical position"


# ---------- new fixes: min font size ----------

def test_available_fixes_offers_min_font_size():
    finding = {
        "slide_number": 4,
        "_skipped": False,
        "min_font_size": {
            "applicable": True, "ok": False,
            "min_required_pt": 9.0, "smallest_pt": 7.0,
            "violations": [{"shape_name": "Body", "size_pt": 7.0}],
        },
        "footer": {}, "action_title": {}, "title_case": {}, "font_family": {},
        "text_length": {},
    }
    fixes = available_fixes_for_slide(finding)
    mfs_fixes = [f for f in fixes if f["id"] == "enforce_min_font_size"]
    assert len(mfs_fixes) == 1
    assert mfs_fixes[0]["min_pt"] == 9.0


# ---------- new fixes: brand font ----------

def test_available_fixes_offers_brand_font_replacement():
    finding = {
        "slide_number": 6,
        "_skipped": False,
        "font_family": {
            "applicable": True, "ok": False,
            "non_brand": [
                {"font": "comic sans ms", "shapes": ["Body"], "fallback_ok": False},
                {"font": "times new roman", "shapes": ["Footer"], "fallback_ok": False},
            ],
        },
        "footer": {}, "action_title": {}, "title_case": {}, "min_font_size": {},
        "text_length": {},
    }
    fixes = available_fixes_for_slide(finding)
    bf_fixes = [f for f in fixes if f["id"] == "enforce_brand_font"]
    assert len(bf_fixes) == 1
    assert bf_fixes[0]["brand_font"] == "ForFuture Sans"


# ---------- new fixes: bullet long paragraph ----------

def test_available_fixes_offers_bullet_split_for_single_long_paragraph():
    finding = {
        "slide_number": 3,
        "_skipped": False,
        "text_length": {
            "ok": False,
            "long_paragraphs": [
                "(CuadroTexto 97, ~4 líneas, 50 palabras) Lorem ipsum dolor sit amet…",
            ],
        },
        "footer": {}, "action_title": {}, "title_case": {}, "min_font_size": {},
        "font_family": {},
    }
    fixes = available_fixes_for_slide(finding)
    bullet_fixes = [f for f in fixes if f["id"] == "bullet_long_paragraph"]
    assert len(bullet_fixes) == 1
    assert bullet_fixes[0]["shape_name"] == "CuadroTexto 97"


def test_available_fixes_does_not_offer_bullet_split_when_multiple_long_paras():
    finding = {
        "slide_number": 3,
        "_skipped": False,
        "text_length": {
            "ok": False,
            "long_paragraphs": [
                "(A, …) one",
                "(B, …) two",
            ],
        },
        "footer": {}, "action_title": {}, "title_case": {}, "min_font_size": {},
        "font_family": {},
    }
    fixes = available_fixes_for_slide(finding)
    assert not any(f["id"] == "bullet_long_paragraph" for f in fixes)
